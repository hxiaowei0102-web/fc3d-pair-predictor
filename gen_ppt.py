#!/usr/bin/env python3
"""生成纯云端部署步骤PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ==================== 颜色方案 ====================
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
TEXT_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_GOLD = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, title, lines, title_color=ACCENT_PURPLE, border_color=None):
    shape = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=border_color or title_color, border_width=Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(12)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = 'Microsoft YaHei'
    
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(4)
    
    return shape

def add_code_block(slide, left, top, width, height, code_lines, title="代码"):
    shape = add_shape(slide, left, top, width, height, fill_color=RGBColor(0x0B, 0x10, 0x1E), border_color=RGBColor(0x33, 0x41, 0x55))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.font.name = 'Consolas'
    
    for line in code_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        p.font.name = 'Consolas'
        p.space_before = Pt(2)
    
    return shape

# ==================== Slide 1: 封面 ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide1, BG_DARK)

# 顶部装饰线
add_shape(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT_RED)

# 标题
add_text_box(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             '福彩3D 两码不组预测系统', font_size=48, color=TEXT_WHITE, bold=True, align=PP_ALIGN.LEFT)

add_text_box(slide1, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             '纯云端部署完全指南', font_size=36, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.LEFT)

add_text_box(slide1, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             'GitHub Pages + GitHub Actions | 零成本 · 全自动 · 手机即用', font_size=20, color=TEXT_GRAY, align=PP_ALIGN.LEFT)

# 特性标签
tags = ['永久免费', '每日自动更新', '手机端优化', '95%准确率', '零运维']
for i, tag in enumerate(tags):
    x = Inches(1 + i * 2.3)
    y = Inches(5)
    card = add_shape(slide1, x, y, Inches(2), Inches(0.6), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
    tf = card.text_frame
    tf.paragraphs[0].text = tag
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = ACCENT_GREEN
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 底部
add_text_box(slide1, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
             '晓炜工作室 · 2026年6月', font_size=14, color=TEXT_MUTED, align=PP_ALIGN.LEFT)

# ==================== Slide 2: 项目架构总览 ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, BG_DARK)

add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             '项目架构总览', font_size=32, color=TEXT_WHITE, bold=True)

add_text_box(slide2, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             '了解整个系统如何运作', font_size=16, color=TEXT_GRAY)

# 架构图 - 4个区块
blocks = [
    ('数据源\ncwl.gov.cn', ACCENT_BLUE, Inches(0.8), Inches(1.8)),
    ('GitHub Actions\n每日21:30运行', ACCENT_GOLD, Inches(4.5), Inches(1.8)),
    ('GitHub Pages\n免费托管网页', ACCENT_GREEN, Inches(8.2), Inches(1.8)),
    ('手机浏览器\n随时随地查看', ACCENT_PURPLE, Inches(11.5), Inches(1.8)),
]

for title, color, x, y in blocks:
    shape = add_shape(slide2, x, y, Inches(2.8), Inches(1.5), fill_color=BG_CARD, border_color=color, border_width=Pt(3))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(20)
    for i, line in enumerate(title.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16) if i == 0 else Pt(12)
        p.font.color.rgb = color if i == 0 else TEXT_GRAY
        p.font.bold = (i == 0)
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER

# 箭头文字
add_text_box(slide2, Inches(3.7), Inches(2.2), Inches(0.8), Inches(0.5), '→', font_size=36, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide2, Inches(7.3), Inches(2.2), Inches(0.8), Inches(0.5), '→', font_size=36, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide2, Inches(10.6), Inches(2.2), Inches(0.8), Inches(0.5), '→', font_size=36, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)

# 详细说明卡片
details = [
    ('📊 V12算法', ['DIGIT 数字投影 (槽1)', 'MOM 数字动量 (槽2)', 'WAKEUP 唤醒检测 (槽3)', '100期回测准确率 95%']),
    ('📁 核心文件', ['auto_predict.py - 自动预测管线', 'predict_v12.py - V12算法引擎', 'build_html.py - HTML生成器', '.github/workflows/predict.yml']),
    ('🔒 安全保障', ['✓ 零未来数据泄漏', '✓ 严格时间序列分离', '✓ 固定规则无过拟合', '✓ 三组算法独立并行']),
]

for i, (title, lines) in enumerate(details):
    add_card(slide2, Inches(0.8 + i * 4.2), Inches(3.8), Inches(3.9), Inches(3.2), title, lines, title_color=ACCENT_GREEN if '安全' in title else ACCENT_PURPLE)

# ==================== Slide 3: 步骤1 - GitHub仓库 ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3, BG_DARK)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             '步骤 1/5', font_size=16, color=ACCENT_GOLD, bold=True)

add_shape(slide3, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05), fill_color=ACCENT_GOLD)

add_text_box(slide3, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             '创建 GitHub 仓库', font_size=32, color=TEXT_WHITE, bold=True)

add_text_box(slide3, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4),
             '免费注册 GitHub 账号，创建一个公开仓库来存放代码', font_size=16, color=TEXT_GRAY)

# 步骤
steps1 = [
    ('1', '注册 GitHub 账号', '访问 github.com → 点击 Sign up → 用邮箱注册\n免费账号即可，无需付费', ACCENT_BLUE),
    ('2', '新建仓库', '点击右上角 + → New repository\n仓库名: fc3d-pair-predictor\n选择 Public（公开）', ACCENT_GREEN),
    ('3', '克隆到本地', 'git clone https://github.com/你的用户名/fc3d-pair-predictor.git\n将项目所有文件复制进去', ACCENT_PURPLE),
]

for i, (num, title, desc, color) in enumerate(steps1):
    y = Inches(2.4 + i * 1.6)
    # 编号圆圈
    circle = slide3.shapes.add_shape(9, Inches(0.8), y + Inches(0.1), Inches(0.6), Inches(0.6))  # oval
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide3, Inches(1.8), y, Inches(3), Inches(0.4), title, font_size=18, color=color, bold=True)
    add_text_box(slide3, Inches(1.8), y + Inches(0.4), Inches(5), Inches(0.8), desc, font_size=13, color=TEXT_GRAY)

# 右侧代码块
add_code_block(slide3, Inches(7.5), Inches(2.4), Inches(5.3), Inches(3.5), [
    '# 项目文件结构',
    '',
    'fc3d-pair-predictor/',
    '├── auto_predict.py    # 自动预测入口',
    '├── predict_v12.py     # V12核心算法',
    '├── build_html.py      # HTML页面生成',
    '├── history.csv        # 历史开奖数据',
    '├── index.html         # 展示页面(生成)',
    '├── prediction_v12.json# 预测结果(生成)',
    '└── .github/',
    '    └── workflows/',
    '        └── predict.yml # 自动任务',
], title="📁 项目文件结构")

# 底部提示
add_text_box(slide3, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             '💡 提示: 仓库必须设为 Public，GitHub Pages 免费服务仅限公开仓库', font_size=14, color=TEXT_MUTED)

# ==================== Slide 4: 步骤2 - GitHub Actions ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide4, BG_DARK)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             '步骤 2/5', font_size=16, color=ACCENT_GOLD, bold=True)
add_shape(slide4, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05), fill_color=ACCENT_GOLD)
add_text_box(slide4, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             '配置自动预测工作流', font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide4, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4),
             '创建 .github/workflows/predict.yml，实现每日自动拉数据→预测→部署', font_size=16, color=TEXT_GRAY)

# 工作流说明
flow_steps = [
    ('⏰ 触发', 'schedule cron\n"30 13 * * *"\n(= BJT 21:30)', ACCENT_BLUE),
    ('📥 拉数据', 'checkout + \nsetup-python\npip install numpy', ACCENT_PURPLE),
    ('🔮 预测', 'python\nauto_predict.py\n↓ 生成JSON+HTML', ACCENT_GOLD),
    ('⬆ 提交', 'git commit\ngit push\n更新到仓库', ACCENT_GREEN),
    ('🚀 部署', 'configure-pages\nupload-artifact\ndeploy-pages', ACCENT_RED),
]

for i, (title, desc, color) in enumerate(flow_steps):
    x = Inches(0.8 + i * 2.5)
    shape = add_shape(slide4, x, Inches(2.3), Inches(2.2), Inches(1.8), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    for line in desc.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER

# 工作流代码
add_code_block(slide4, Inches(0.8), Inches(4.4), Inches(6), Inches(2.8), [
    'name: 福彩3D每日预测',
    'on:',
    '  schedule:',
    '    - cron: "30 13 * * *"  # BJT 21:30',
    '  workflow_dispatch:',  # 手动触发
    '',
    'jobs:',
    '  predict-and-deploy:',
    '    runs-on: ubuntu-latest',
    '    permissions:',
    '      contents: write',
    '      pages: write',
    '    steps:',
    '      - uses: actions/checkout@v4',
    '      - uses: actions/setup-python@v5',
    '        with: python-version: "3.11"',
    '      - run: pip install numpy',
    '      - run: python auto_predict.py',
    '      - run: |',
    '          git config user.name "bot"',
    '          git add index.html prediction_v12.json',
    '          git commit -m "🔮 auto" || true',
    '          git push',
], title="📄 .github/workflows/predict.yml")

# 右侧说明
add_card(slide4, Inches(7.2), Inches(4.4), Inches(5.5), Inches(2.8), 
    '⚙️ 关键配置说明',
    [
        'cron "30 13 * * *" = UTC 13:30 = BJT 21:30',
        'workflow_dispatch = 可随时手动触发',
        'permissions: write = 允许自动提交代码',
        'ubuntu-latest = GitHub提供的免费虚拟机',
        'numpy = 算法计算所需的唯一依赖',
        '部署到 Pages 无需额外配置域名或服务器',
    ],
    title_color=ACCENT_GOLD,
    border_color=ACCENT_GOLD
)

# ==================== Slide 5: 步骤3 - GitHub Pages ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide5, BG_DARK)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             '步骤 3/5', font_size=16, color=ACCENT_GOLD, bold=True)
add_shape(slide5, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05), fill_color=ACCENT_GOLD)
add_text_box(slide5, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             '启用 GitHub Pages', font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide5, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4),
             '将仓库自动发布为网站，获得永久免费的 HTTPS 域名', font_size=16, color=TEXT_GRAY)

# Pages设置步骤
pages_steps = [
    ('1️⃣', '进入仓库 Settings', '打开 GitHub 仓库 → 点击顶部 Settings 标签页\n在左侧菜单找到 Pages 选项', ACCENT_BLUE),
    ('2️⃣', '选择部署源', 'Source: Deploy from a branch\nBranch: main / root\n点击 Save', ACCENT_GREEN),
    ('3️⃣', '等待部署', '首次部署约 1-2 分钟\n刷新页面即可看到绿色提示:\n"Your site is live at..."', ACCENT_GOLD),
    ('4️⃣', '获取网址', 'https://你的用户名.github.io/仓库名/\n这就是你的永久访问地址\n支持 HTTPS 加密', ACCENT_PURPLE),
]

for i, (num, title, desc, color) in enumerate(pages_steps):
    y = Inches(2.3 + i * 1.3)
    shape = add_shape(slide5, Inches(0.8), y, Inches(5.5), Inches(1.1), fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    # 编号
    add_text_box(slide5, Inches(1.0), y + Inches(0.05), Inches(0.6), Inches(0.5), num, font_size=24, color=color, bold=True)
    add_text_box(slide5, Inches(1.8), y + Inches(0.05), Inches(2), Inches(0.4), title, font_size=16, color=color, bold=True)
    add_text_box(slide5, Inches(1.8), y + Inches(0.45), Inches(4.3), Inches(0.6), desc, font_size=12, color=TEXT_GRAY)

# 右侧实际地址
add_card(slide5, Inches(7), Inches(2.3), Inches(5.5), Inches(3.2),
    '🌐 实际部署地址',
    [
        '项目网址:',
        'hxiaowei0102-web.github.io/fc3d-pair-predictor/',
        '',
        '✓ 永久免费，无需续费',
        '✓ 自动 HTTPS 加密',
        '✓ 全球 CDN 加速',
        '✓ 支持手机直接访问',
        '✓ 每天 21:30 自动更新内容',
    ],
    title_color=ACCENT_GREEN,
    border_color=ACCENT_GREEN
)

# ==================== Slide 6: 步骤4 - 手机使用 ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide6, BG_DARK)

add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             '步骤 4/5', font_size=16, color=ACCENT_GOLD, bold=True)
add_shape(slide6, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05), fill_color=ACCENT_GOLD)
add_text_box(slide6, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             '手机端使用指南', font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide6, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4),
             '只需手机浏览器，无需安装任何 App', font_size=16, color=TEXT_GRAY)

# iOS 和 Android 指南
add_card(slide6, Inches(0.8), Inches(2.3), Inches(5.8), Inches(3), 
    '🍎 iPhone / iOS', 
    [
        '1. 打开 Safari 浏览器',
        '2. 输入网址打开页面',
        '3. 点击底部 分享按钮 (↑)',
        '4. 选择「添加到主屏幕」',
        '5. 命名后点击「添加」',
        '6. 桌面出现 App 图标，点击即用',
        '',
        '✨ 已配置: apple-mobile-web-app-capable',
        '全屏体验，隐藏浏览器工具栏',
    ],
    title_color=ACCENT_BLUE,
    border_color=ACCENT_BLUE
)

add_card(slide6, Inches(7), Inches(2.3), Inches(5.8), Inches(3), 
    '🤖 Android / 安卓', 
    [
        '1. 打开 Chrome 浏览器',
        '2. 输入网址打开页面',
        '3. 点击右上角 三个点 (⋮)',
        '4. 选择「添加到主屏幕」',
        '5. 确认名称后点击「添加」',
        '6. 桌面出现快捷方式，点击即用',
        '',
        '✨ 页面已针对 640px 移动端优化',
        '每 5 分钟自动刷新数据',
    ],
    title_color=ACCENT_GREEN,
    border_color=ACCENT_GREEN
)

# 底部特性
add_card(slide6, Inches(0.8), Inches(5.6), Inches(11.9), Inches(1.5), 
    '📱 页面核心功能', 
    [
        '① 首行大字显示期号 | ② 三组独立算法预测数字 (DIGIT/MOM/WAKEUP) | ③ 95%准确率统计卡片',
        '④ 算法详细介绍区域 | ⑤ 100期回测明细表格 (由近到远) | ⑥ 底部显示更新时间和版本历程',
    ],
    title_color=ACCENT_PURPLE,
    border_color=ACCENT_PURPLE
)

# ==================== Slide 7: 步骤5 - 费用与维护 ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide7, BG_DARK)

add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             '步骤 5/5', font_size=16, color=ACCENT_GOLD, bold=True)
add_shape(slide7, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05), fill_color=ACCENT_GOLD)
add_text_box(slide7, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             '费用与日常维护', font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide7, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4),
             '全部免费，零运维负担', font_size=16, color=TEXT_GRAY)

# 费用表格
costs = [
    ('GitHub Pages', '网页托管', '永久免费', '无流量限制', ACCENT_GREEN),
    ('GitHub Actions', '自动执行任务', '2000分钟/月免费', '我们只用~30分钟/月', ACCENT_GREEN),
    ('域名', 'github.io 子域名', '永久免费', 'HTTPS 自动配置', ACCENT_GREEN),
    ('服务器', '不需要', '¥0', 'GitHub 提供虚拟机', ACCENT_GREEN),
    ('数据库', '不需要', '¥0', 'CSV 文件存储', ACCENT_GREEN),
    ('总计', '', '永久 ¥0/月', '', ACCENT_GOLD),
]

# 表头
for j, (h, w) in enumerate([('服务项目', 2.5), ('用途', 2.5), ('费用', 2.5), ('说明', 3), ('状态', 1.5)]):
    x = Inches(1.2 + sum([2.5, 2.5, 2.5, 3, 1.5][:j]))-Inches(2.5)+Inches(j*2.7)
    add_text_box(slide7, Inches(1.2 + j * 2.6), Inches(2.3), Inches(2.4), Inches(0.5), h, font_size=15, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)

for i, (svc, purpose, cost, note, color) in enumerate(costs):
    y = Inches(2.8 + i * 0.55)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x16, 0x21, 0x33)
    if i == len(costs) - 1:
        bg = RGBColor(0x2D, 0x1F, 0x0E)
    
    for j, txt in enumerate([svc, purpose, cost, note]):
        shape = add_shape(slide7, Inches(1.2 + j * 2.6), y, Inches(2.4), Inches(0.5), fill_color=bg)
        tf = shape.text_frame
        tf.paragraphs[0].text = txt
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = color if i == len(costs)-1 else TEXT_WHITE
        tf.paragraphs[0].font.bold = (i == len(costs)-1 or j == 0)
        tf.paragraphs[0].font.name = 'Microsoft YaHei'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 状态
    shape = add_shape(slide7, Inches(12), y, Inches(1), Inches(0.5), fill_color=bg)
    tf = shape.text_frame
    tf.paragraphs[0].text = '✓ 免费' if cost != '' else ''
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = ACCENT_GREEN
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 维护说明
add_card(slide7, Inches(0.8), Inches(5.8), Inches(5.8), Inches(1.3), 
    '🔄 日常维护 (几乎为零)',
    [
        '✓ 系统全自动运行，无需人工干预',
        '✓ 每天 21:30 BJT 自动拉取最新开奖数据',
        '✓ 自动执行 V12 算法预测 + 100期回测',
        '✓ 自动生成 HTML 页面并部署到 GitHub Pages',
    ],
    title_color=ACCENT_GREEN,
    border_color=ACCENT_GREEN
)

add_card(slide7, Inches(7), Inches(5.8), Inches(5.8), Inches(1.3), 
    '⚠️ 可能需要关注的情况',
    [
        '• 如果 GitHub Action 运行失败（极少发生）',
        '  → GitHub 会发邮件通知',
        '• 如果数据源 API 故障',
        '  → 系统自动回退使用本地历史数据',
    ],
    title_color=ACCENT_GOLD,
    border_color=ACCENT_GOLD
)

# ==================== Slide 8: 总结 ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide8, BG_DARK)

add_text_box(slide8, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             '总结', font_size=32, color=TEXT_WHITE, bold=True)
add_shape(slide8, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), fill_color=ACCENT_GOLD)

# 三步总结
summary = [
    ('📂 代码', '推送到 GitHub', 'git add . && git commit && git push', ACCENT_BLUE),
    ('⚙️ 自动', 'GitHub Actions', '每日 21:30 自动执行', ACCENT_GOLD),
    ('🌐 访问', '手机打开网址', 'github.io/fc3d-pair-predictor', ACCENT_GREEN),
]

for i, (icon, title, desc, color) in enumerate(summary):
    x = Inches(0.8 + i * 4.2)
    shape = add_shape(slide8, x, Inches(1.6), Inches(3.8), Inches(2.2), fill_color=BG_CARD, border_color=color, border_width=Pt(3))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(20)
    
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)

# 关键数据
add_text_box(slide8, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
             '当前运行数据', font_size=20, color=TEXT_WHITE, bold=True)

data_cards = [
    ('95.0%', '100期准确率', ACCENT_RED),
    ('1005期', '训练数据量', ACCENT_BLUE),
    ('V12', '算法版本', ACCENT_GOLD),
    ('21:30', '每日更新时间', ACCENT_GREEN),
    ('¥0', '月度费用', ACCENT_PURPLE),
    ('3组', '独立算法', ACCENT_RED),
]

for i, (val, label, color) in enumerate(data_cards):
    x = Inches(0.8 + i * 2.1)
    shape = add_shape(slide8, x, Inches(4.8), Inches(1.9), Inches(1.2), fill_color=BG_CARD, border_color=color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(30)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_GRAY
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

# 底部网址
shape = add_shape(slide8, Inches(2), Inches(6.3), Inches(9.3), Inches(0.8), fill_color=RGBColor(0x2D, 0x1F, 0x0E), border_color=ACCENT_GOLD, border_width=Pt(2))
tf = shape.text_frame
tf.paragraphs[0].text = '🔗 hxiaowei0102-web.github.io/fc3d-pair-predictor/'
tf.paragraphs[0].font.size = Pt(22)
tf.paragraphs[0].font.color.rgb = ACCENT_GOLD
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Consolas'
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ==================== 保存 ====================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '福彩3D纯云端部署指南.pptx')
prs.save(output_path)
print(f"PPT saved: {output_path}")
